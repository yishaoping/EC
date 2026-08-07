from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR


OUT = "/data1/gzh/EC/BOOM_Rocket_Barrier_Rollback.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

FONT = "Noto Sans CJK SC"
MONO = "Noto Sans Mono CJK SC"

BG = RGBColor(13, 20, 29)
PANEL = RGBColor(24, 34, 46)
PANEL2 = RGBColor(31, 43, 57)
WHITE = RGBColor(241, 245, 248)
MUTED = RGBColor(165, 181, 193)
CYAN = RGBColor(69, 201, 214)
CYAN_D = RGBColor(24, 87, 102)
ORANGE = RGBColor(242, 157, 70)
ORANGE_D = RGBColor(113, 66, 31)
RED = RGBColor(241, 104, 101)
RED_D = RGBColor(105, 43, 48)
YELLOW = RGBColor(239, 203, 90)
YELLOW_D = RGBColor(104, 86, 26)
GREEN = RGBColor(96, 202, 145)
GREEN_D = RGBColor(33, 84, 64)
GRID = RGBColor(53, 69, 84)


def rgb(c):
    return c


def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def text_box(slide, x, y, w, h, text, size=16, color=WHITE, bold=False,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, font=FONT,
             margin=0.06, italic=False):
    tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tx


def rich_lines(slide, x, y, w, h, lines, size=13, color=WHITE, bullet=False,
               line_spacing=1.15, margin=0.08):
    tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    for i, item in enumerate(lines):
        if isinstance(item, tuple):
            txt, col, b = item
        else:
            txt, col, b = item, color, False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        p.line_spacing = line_spacing
        if bullet:
            p.text = "• " + txt
        else:
            p.text = txt
        for r in p.runs:
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.color.rgb = col
            r.font.bold = b
    return tx


def box(slide, x, y, w, h, fill=PANEL, line=GRID, radius=0.10, lw=1.0):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(lw)
    return shape


def line(slide, x1, y1, x2, y2, color=GRID, width=1.2, dash=None):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(width)
    if dash:
        ln.line.dash_style = dash
    return ln


def arrow(slide, x1, y1, x2, y2, color=CYAN, width=1.8):
    ln = line(slide, x1, y1, x2, y2, color, width)
    ln.line.end_arrowhead = True
    return ln


def tag(slide, x, y, w, label, fill, color=BG, size=10):
    s = box(slide, x, y, w, 0.27, fill, fill, 0.10, 0.5)
    text_box(slide, x, y+0.01, w, 0.22, label, size, color, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, margin=0.01)
    return s


def title(slide, no, kicker, heading, statement, source=None):
    text_box(slide, 0.54, 0.25, 0.6, 0.28, f"{no:02d}", 11, CYAN, True, font=MONO, margin=0)
    text_box(slide, 1.15, 0.23, 3.4, 0.28, kicker.upper(), 10, MUTED, True, font=MONO, margin=0)
    text_box(slide, 0.54, 0.57, 12.2, 0.48, heading, 27, WHITE, True, margin=0)
    # conclusion ribbon
    box(slide, 0.54, 1.14, 12.25, 0.47, CYAN_D, CYAN_D, 0.08, 0.5)
    text_box(slide, 0.73, 1.20, 11.85, 0.33, "核心观点  " + statement, 13, WHITE, True, valign=MSO_ANCHOR.MIDDLE, margin=0)
    if source:
        text_box(slide, 5.05, 0.25, 7.70, 0.18, "REF · " + source, 6.6, MUTED, font=MONO, align=PP_ALIGN.RIGHT, margin=0)


def footer(slide, section):
    line(slide, 0.55, 7.01, 12.78, 7.01, GRID, 0.7)
    text_box(slide, 0.56, 7.08, 4.0, 0.18, "Chipyard · BOOM ↔ Rocket · Error recovery", 7.5, MUTED, font=MONO, margin=0)
    text_box(slide, 10.5, 7.08, 2.25, 0.18, section, 7.5, MUTED, font=MONO, align=PP_ALIGN.RIGHT, margin=0)


def node(slide, x, y, w, h, name, sub=None, fill=PANEL2, edge=GRID, accent=None, name_size=13):
    box(slide, x, y, w, h, fill, edge, 0.10, 1.1)
    if accent:
        slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.07), Inches(h)).fill.solid()
        ac = slide.shapes[-1]
        ac.fill.fore_color.rgb = accent
        ac.line.fill.background()
    text_box(slide, x+0.14, y+0.13, w-0.28, 0.28, name, name_size, WHITE, True, valign=MSO_ANCHOR.MIDDLE, margin=0)
    if sub:
        text_box(slide, x+0.14, y+0.48, w-0.28, h-0.54, sub, 9.5, MUTED, margin=0)


def bullet_card(slide, x, y, w, h, head, items, accent=CYAN, fill=PANEL):
    box(slide, x, y, w, h, fill, GRID, 0.08, 0.8)
    text_box(slide, x+0.15, y+0.12, w-0.3, 0.28, head, 13, accent, True, margin=0)
    rich_lines(slide, x+0.14, y+0.47, w-0.28, h-0.58, items, 10.5, WHITE, bullet=True, margin=0.0)


def add_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    return slide


# Slide 1
s = add_slide()
text_box(s, 0.62, 0.50, 4.5, 0.3, "CHIPYARD RESEARCH BRIEF · 2026", 10, CYAN, True, font=MONO, margin=0)
text_box(s, 0.62, 1.08, 6.20, 1.52, "面向 BOOM–Rocket 异构核\n协同校验系统的错误传播分析\n与屏障回滚机制设计", 24, WHITE, True, margin=0)
text_box(s, 0.65, 2.72, 6.25, 0.4, "把“检测到错误”推进为“错误不逃逸、状态可恢复”", 15, MUTED, margin=0)

# hero architecture diagram
box(s, 7.2, 1.22, 5.42, 4.95, PANEL, GRID, 0.10, 0.8)
text_box(s, 7.52, 1.47, 4.8, 0.26, "异构协同校验闭环", 13, WHITE, True, margin=0)
node(s, 7.65, 2.05, 1.65, 1.12, "BOOM", "1 × high-performance\n200 MHz · hart 0", PANEL2, CYAN, CYAN, 15)
node(s, 10.42, 2.05, 1.72, 1.12, "Rocket", "4 × checker\n100 MHz · hart 1–4", PANEL2, CYAN, GREEN, 15)
arrow(s, 9.32, 2.58, 10.35, 2.58, CYAN, 2.5)
text_box(s, 9.20, 2.18, 1.25, 0.24, "校验 packet", 9.5, CYAN, True, align=PP_ALIGN.CENTER, margin=0)
node(s, 8.98, 3.55, 2.10, 0.78, "GH_BUF → GHM/CDC", "136 bit × 2 packet · depth 256", PANEL2, GRID, YELLOW, 11)
arrow(s, 8.48, 3.18, 9.60, 3.50, CYAN, 1.8)
arrow(s, 10.93, 3.18, 10.53, 3.50, GREEN, 1.8)
node(s, 7.65, 4.88, 1.65, 0.76, "L1D / L2", "cacheable state", PANEL2, GRID, ORANGE, 12)
node(s, 10.42, 4.88, 1.72, 0.76, "MMIO / DRAM", "external state", PANEL2, GRID, RED, 12)
arrow(s, 8.47, 5.25, 10.34, 5.25, ORANGE, 1.9)
text_box(s, 7.60, 5.76, 4.85, 0.25, "Barrier gates irreversible effects; rollback repairs reversible effects", 9.3, MUTED, align=PP_ALIGN.CENTER, margin=0)

bullet_card(s, 0.65, 3.42, 2.05, 1.82, "问题", ["BOOM 先执行，Rocket 后确认", "副作用可能先于校验发生"], RED)
bullet_card(s, 2.90, 3.42, 2.05, 1.82, "目标", ["隔离未确认状态", "错误后回到一致点"], CYAN)
bullet_card(s, 5.15, 3.42, 1.72, 1.82, "方法", ["路径分层", "难度分级"], ORANGE)
text_box(s, 0.66, 6.18, 6.2, 0.42, "核心抽象：Propagation Path × Rollback Difficulty", 15, WHITE, True, margin=0)
footer(s, "01 / BACKGROUND")

# Slide 2
s = add_slide()
title(s, 2, "Problem framing", "为什么需要 Barrier + Rollback？", "屏障阻断不可逆传播；回滚撤销已经发生的可逆修改。", "Method/EC.typ；Method/boom_dataflow_analysis.typ:51–65")

# left event timeline
text_box(s, 0.70, 1.83, 5.7, 0.25, "错误窗口：执行先于确认", 13, WHITE, True, margin=0)
line(s, 1.0, 3.02, 6.10, 3.02, GRID, 1.4)
for x, lab, col in [(1.10, "执行", CYAN), (2.45, "写入 L1D", ORANGE), (3.95, "传播", RED), (5.52, "Rocket 确认", GREEN)]:
    slide = s
    slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(2.84), Inches(0.36), Inches(0.36)).fill.solid()
    dot = slide.shapes[-1]
    dot.fill.fore_color.rgb = col; dot.line.color.rgb = col
    text_box(s, x-0.22, 3.34, 1.0, 0.24, lab, 10, col, True, align=PP_ALIGN.CENTER, margin=0)
text_box(s, 0.93, 2.22, 1.25, 0.42, "BOOM\nissue/execute", 10, WHITE, True, align=PP_ALIGN.CENTER, margin=0)
text_box(s, 2.18, 2.22, 1.55, 0.42, "dirty line\n形成", 10, ORANGE, True, align=PP_ALIGN.CENTER, margin=0)
text_box(s, 3.66, 2.22, 1.65, 0.42, "L2 / 外设\n可见", 10, RED, True, align=PP_ALIGN.CENTER, margin=0)
text_box(s, 5.18, 2.22, 1.52, 0.42, "校验结果\n到达", 10, GREEN, True, align=PP_ALIGN.CENTER, margin=0)

box(s, 0.75, 4.15, 2.55, 1.55, RED_D, RED, 0.08, 1.0)
text_box(s, 0.94, 4.35, 2.16, 0.28, "没有 Barrier", 14, RED, True, align=PP_ALIGN.CENTER, margin=0)
rich_lines(s, 0.95, 4.78, 2.13, 0.72, ["MMIO Put/Atomic 立即生效", "DRAM 写回后难以追回"], 10.5, WHITE, bullet=True, margin=0)

box(s, 3.55, 4.15, 2.55, 1.55, ORANGE_D, ORANGE, 0.08, 1.0)
text_box(s, 3.74, 4.35, 2.16, 0.28, "没有 Rollback", 14, ORANGE, True, align=PP_ALIGN.CENTER, margin=0)
rich_lines(s, 3.75, 4.78, 2.12, 0.72, ["L1D dirty line 保留错误值", "CSR / reservation 无法重建"], 10.5, WHITE, bullet=True, margin=0)

# right principle panel
box(s, 6.65, 1.84, 5.92, 4.86, PANEL, GRID, 0.10, 0.8)
text_box(s, 6.98, 2.12, 5.1, 0.35, "设计问题应改写为 3 个追问", 15, WHITE, True, margin=0)
questions = [("01", "改写了什么？", "store / SC / AMO / CSR / reservation", RED),
             ("02", "沿哪条路径？", "L1D → L2 → DRAM  或  IOMSHR → 外设", ORANGE),
             ("03", "在哪里止损？", "Barrier before commit  +  checkpoint / undo", CYAN)]
for i, (n, q, d, col) in enumerate(questions):
    yy = 2.75 + i*1.08
    text_box(s, 7.05, yy, 0.45, 0.28, n, 11, col, True, font=MONO, margin=0)
    text_box(s, 7.68, yy-0.03, 2.2, 0.30, q, 13, col, True, margin=0)
    text_box(s, 7.68, yy+0.35, 4.35, 0.35, d, 10, MUTED, margin=0)
    if i < 2: line(s, 7.05, yy+0.82, 12.15, yy+0.82, GRID, 0.8)
text_box(s, 7.05, 6.15, 5.0, 0.28, "正确性目标：error detected  ≠  error escaped", 12, WHITE, True, margin=0)
footer(s, "02 / WHY RECOVERY")

# Slide 3
s = add_slide()
title(s, 3, "Level 0", "BOOM 执行过程中：哪些状态对外可见？", "重点不是所有输出，而是能改变系统可观察状态的操作。", "Method/boom_dataflow_analysis.typ:164–190")

# central core
node(s, 5.36, 3.03, 2.62, 1.12, "BOOM core", "ROB / LSU / CSR / L1D", PANEL2, CYAN, CYAN, 18)
# left outputs
for yy, name, sub, col in [(1.90, "校验 packet", "执行结果 · load/store 信息", CYAN),
                           (3.35, "Memory side effect", "Store · Load · SC · LR · AMO", ORANGE),
                           (4.80, "Architectural state", "CSR · privilege · register state", YELLOW)]:
    node(s, 0.78, yy, 3.55, 0.90, name, sub, PANEL, col, col, 13)
    arrow(s, 4.35, yy+0.45, 5.25, 3.55, col, 1.5)
# right outcomes
for yy, name, sub, col in [(2.15, "Rocket checker", "比较 / 重放 / 通过或拒绝", GREEN),
                           (3.70, "Memory hierarchy", "L1D → L2 → DRAM", ORANGE),
                           (5.10, "外设 / MMIO", "不可缓存、不可逆副作用", RED)]:
    node(s, 9.02, yy, 3.40, 0.90, name, sub, PANEL, col, col, 13)
    arrow(s, 8.08, 3.55, 8.92, yy+0.45, col, 1.5)
# classify strip
tag(s, 0.78, 6.22, 1.12, "校验路径", CYAN, BG, 9)
text_box(s, 2.00, 6.23, 2.2, 0.25, "可比较，但不等于已提交", 9.5, MUTED, margin=0)
tag(s, 4.70, 6.22, 1.25, "缓存路径", ORANGE, BG, 9)
text_box(s, 6.03, 6.23, 2.2, 0.25, "先行可隐藏，外发必须受控", 9.5, MUTED, margin=0)
tag(s, 8.73, 6.22, 1.18, "外设路径", RED, BG, 9)
text_box(s, 9.98, 6.23, 2.4, 0.25, "一旦到达，通常不可回滚", 9.5, MUTED, margin=0)
footer(s, "03 / OBSERVABLE STATE")

# Slide 4
s = add_slide()
title(s, 4, "Propagation model", "错误传播路径：从 BOOM 到外部世界的四层模型", "每跨过一层边界，恢复成本上升；外设是最外层、最难逆转的屏障。", "Method/boom_dataflow_analysis.typ:132–162")

# layered flow diagram
layers = [
    ("L0", "BOOM outputs", "待校验 packet · store/load/SC/LR/AMO · CSR", CYAN),
    ("L1", "Tile-local state", "L1D data/metadata · MSHR · IOMSHR", ORANGE),
    ("L2", "Shared boundary", "L2 coherence · SBUS/CBUS/PBUS", ORANGE),
    ("L3", "System outside", "DRAM · UART · CLINT/PLIC · external device", RED),
]
for i, (lv, name, desc, col) in enumerate(layers):
    yy = 1.87 + i*1.12
    box(s, 0.82, yy, 11.68, 0.82, PANEL if i < 3 else RED_D, col, 0.08, 1.1)
    tag(s, 1.02, yy+0.27, 0.58, lv, col, BG, 10)
    text_box(s, 1.85, yy+0.13, 2.25, 0.26, name, 13, col, True, margin=0)
    text_box(s, 4.18, yy+0.15, 7.86, 0.30, desc, 11, WHITE, margin=0)
    if i < 3:
        arrow(s, 6.65, yy+0.83, 6.65, yy+1.08, col, 1.4)
        text_box(s, 8.85, yy+0.84, 2.8, 0.20, "边界成本 ↑", 8.5, MUTED, font=MONO, margin=0)
# gate markers
box(s, 11.02, 2.04, 1.15, 0.48, CYAN_D, CYAN, 0.07, 0.8)
text_box(s, 11.07, 2.14, 1.05, 0.2, "check", 10, CYAN, True, align=PP_ALIGN.CENTER, margin=0)
box(s, 11.02, 4.28, 1.15, 0.48, RED_D, RED, 0.07, 0.8)
text_box(s, 11.06, 4.38, 1.07, 0.2, "barrier", 10, RED, True, align=PP_ALIGN.CENTER, margin=0)

box(s, 0.82, 6.55, 11.68, 0.34, PANEL2, GRID, 0.06, 0.5)
text_box(s, 1.02, 6.60, 11.2, 0.18, "关键判断点：dmem.req.fire → data-array write → C Release/ProbeAckData → L2 writeback → AXI W/B", 9.3, WHITE, font=MONO, align=PP_ALIGN.CENTER, margin=0)
footer(s, "04 / PROPAGATION PATH")

# Slide 5
s = add_slide()
title(s, 5, "Memory hierarchy", "Memory Hierarchy 中：同一条 store 如何越过屏障？", "cacheable 写入可以暂存在层次内部；uncacheable 写入在 A 通道前必须 gate。", "Method/boom_dataflow_analysis.typ:287–412；Method/boom_traffic_statistics.typ")

# two branch diagrams
text_box(s, 0.78, 1.82, 5.7, 0.25, "A. Cacheable：可隐藏，但不能任其外发", 13, ORANGE, True, margin=0)
flow = [("STQ / LSU", 0.82, 2.34, 1.42, 0.64, CYAN), ("L1D\ndata array", 2.58, 2.34, 1.45, 0.64, ORANGE), ("C channel\nRelease/ProbeAckData", 4.35, 2.34, 1.65, 0.64, ORANGE), ("L2\nshared", 6.33, 2.34, 1.22, 0.64, ORANGE), ("DRAM", 7.90, 2.34, 1.10, 0.64, RED)]
for name, x, y, w, h, col in flow:
    node(s, x, y, w, h, name, None, PANEL, col, col, 10.5)
for x1, x2, col in [(2.27,2.51,CYAN),(4.05,4.28,ORANGE),(6.00,6.27,ORANGE),(7.60,7.84,RED)]: arrow(s,x1,2.66,x2,2.66,col,1.8)
box(s, 2.36, 3.44, 5.50, 0.62, ORANGE_D, ORANGE, 0.07, 0.8)
text_box(s, 2.56, 3.61, 5.10, 0.24, "未确认 dirty line：允许留在 L1D，禁止 Release / Probe 外发", 11, WHITE, True, align=PP_ALIGN.CENTER, margin=0)
text_box(s, 0.82, 4.25, 8.3, 0.28, "控制点：line timestamp / epoch + byte mask + eviction / Probe gate", 10.5, MUTED, font=MONO, margin=0)

text_box(s, 9.37, 1.82, 3.1, 0.25, "B. Uncacheable：外设前的硬屏障", 13, RED, True, margin=0)
uflow = [("STQ / LSU", 9.37, 2.34, 1.12, CYAN), ("IOMSHR", 10.73, 2.34, 1.08, RED), ("A channel\nPut / Atomic", 9.90, 3.45, 1.45, RED), ("外设", 11.65, 3.45, 0.82, RED)]
node(s, 9.37, 2.34, 1.12, 0.64, "STQ / LSU", None, PANEL, CYAN, CYAN, 10.5)
node(s, 10.73, 2.34, 1.08, 0.64, "IOMSHR", None, PANEL, RED, RED, 10.5)
node(s, 9.90, 3.45, 1.45, 0.64, "A channel\nPut / Atomic", None, PANEL, RED, RED, 10.5)
node(s, 11.65, 3.45, 0.82, 0.64, "外设", "UART / PLIC", PANEL, RED, RED, 10.5)
arrow(s,10.51,2.66,10.65,2.66,RED,1.8); arrow(s,11.27,2.99,10.95,3.40,RED,1.8); arrow(s,11.35,3.77,11.57,3.77,RED,1.8)
box(s, 9.55, 4.56, 2.77, 0.90, RED_D, RED, 0.07, 1.0)
text_box(s, 9.70, 4.71, 2.48, 0.28, "Barrier before commit", 12, RED, True, align=PP_ALIGN.CENTER, margin=0)
text_box(s, 9.72, 5.08, 2.45, 0.25, "等待 Rocket 通过后再发 A", 9.8, WHITE, align=PP_ALIGN.CENTER, margin=0)
text_box(s, 9.44, 5.82, 2.95, 0.50, "MMIO 一旦被 manager 接收，\ncore rollback 无法撤回设备副作用", 10.2, MUTED, align=PP_ALIGN.CENTER, margin=0)

box(s, 0.82, 5.60, 8.25, 0.82, PANEL2, GRID, 0.07, 0.7)
text_box(s, 1.00, 5.79, 7.88, 0.36, "统计口径：`dmem_req_fire` = 指令请求；`dataWriteArb.fire` = L1D 真正修改；`tl_out.c.fire` = cache line 开始外发", 10, WHITE, margin=0)
footer(s, "05 / MEMORY HIERARCHY")

# Slide 6
s = add_slide()
title(s, 6, "Rollback difficulty", "五级恢复难度：从不可提前写出到无影响操作", "恢复难度由“外部可见性 + 可逆性 + 状态耦合度”共同决定。", "Method/EC.typ；Method/boom_dataflow_analysis.typ:247–279")

levels = [
    ("L1", "不可提前写出", "uncacheable Store / AMO", "必须等校验", RED),
    ("L2", "额外状态修改", "uncacheable Load · LR/SC reservation", "备份返回值 / reservation", YELLOW),
    ("L3", "Cache 系统状态", "cacheable Store · dirty line · metadata", "line / epoch / eviction", ORANGE),
    ("L4", "架构状态恢复", "CSR · register · privilege · PC", "checkpoint 快照", CYAN),
    ("L5", "无外部影响", "cacheable Load", "replay / re-check", GREEN),
]
for i, (lv, head, ops, strat, col) in enumerate(levels):
    x = 0.78 + i*2.47
    box(s, x, 2.06, 2.10, 3.86, PANEL if i not in (0,4) else (RED_D if i==0 else GREEN_D), col, 0.08, 1.1)
    tag(s, x+0.18, 2.26, 0.53, lv, col, BG, 10)
    text_box(s, x+0.18, 2.72, 1.72, 0.56, head, 14, col, True, margin=0)
    line(s, x+0.18, 3.42, x+1.92, 3.42, GRID, 0.8)
    text_box(s, x+0.18, 3.70, 1.74, 0.86, ops, 11, WHITE, True, margin=0)
    text_box(s, x+0.18, 4.84, 1.76, 0.48, strat, 10, MUTED, margin=0)
    # difficulty bar
    for j in range(i+1):
        box(s, x+0.18+j*0.29, 5.55, 0.22, 0.14, col, col, 0.02, 0.3)
    text_box(s, x+0.18, 5.73, 1.74, 0.17, "恢复成本", 8.2, MUTED, font=MONO, margin=0)

text_box(s, 0.85, 6.30, 11.6, 0.30, "边界提醒：普通 load 本身不写外部状态；但 read-clear / FIFO pop / 动态 CSR 属于不可重复输入，需单独记录。", 10.2, YELLOW, True, margin=0)
footer(s, "06 / DIFFICULTY TAXONOMY")

# Slide 7
s = add_slide()
title(s, 7, "Recovery policy", "不同类别对应什么恢复策略？", "策略不是“一刀切暂停”：对不可逆路径加屏障，对可隐藏路径加版本与撤销能力。", "Method/boom_dataflow_analysis.typ:417–455")

headers = [(0.78, 2.48, "状态类别", CYAN), (3.36, 3.15, "允许提前做什么", ORANGE), (6.60, 3.05, "必须记录什么", YELLOW), (9.78, 2.75, "错误时怎么恢复", GREEN)]
for x,w,h,col in headers:
    box(s, x, 1.92, w, 0.48, PANEL2, col, 0.06, 0.8)
    text_box(s, x+0.12, 2.03, w-0.24, 0.22, h, 11, col, True, align=PP_ALIGN.CENTER, margin=0)
rows = [
    ("uncacheable Store / AMO", "不发 A 通道请求", "地址、数据、opcode、顺序", "Barrier gate；校验通过再 commit", RED),
    ("uncacheable Load / MMIO read", "可执行一次并缓存返回", "返回值、事件序号、side-effect 标记", "回放复用，不重复读设备", YELLOW),
    ("cacheable Store / SC / AMO", "可进入 L1D 私有状态", "old value、byte mask、epoch、line metadata", "undo log + line restore；禁止外发", ORANGE),
    ("CSR / register / privilege", "在 speculative window 内推进", "architectural checkpoint", "恢复 PC + CSR + register state", CYAN),
    ("cacheable Load", "自由执行 / 重放", "可选的依赖与顺序信息", "重新执行或重新校验", GREEN),
]
for i, row in enumerate(rows):
    yy = 2.52 + i*0.82
    fill = PANEL if i%2==0 else PANEL2
    for (x,w), val in zip([(0.78,2.48),(3.36,3.15),(6.60,3.05),(9.78,2.75)], row[:4]):
        box(s, x, yy, w, 0.70, fill, GRID, 0.03, 0.5)
        text_box(s, x+0.11, yy+0.12, w-0.22, 0.45, val, 9.5, WHITE if x != 0.78 else row[4], bold=(x==0.78), margin=0)
text_box(s, 0.83, 6.72, 11.5, 0.22, "统一判据：凡是可能到达不可恢复区域的修改，必须先验证；凡是允许提前执行的修改，必须有 Undo。", 10.5, WHITE, True, align=PP_ALIGN.CENTER, margin=0)
footer(s, "07 / RECOVERY POLICY")

# Slide 8
s = add_slide()
title(s, 8, "End-to-end framework", "整体 BOOM–Rocket 协同校验与恢复框架", "校验路径决定“何时放行”，状态日志决定“错误后能否回到一致点”。", "Method/boom_dataflow_analysis.typ:312–342；Method/ParaMedic配置对比与瓶颈分析.typ")

# top path
node(s, 0.80, 1.93, 1.48, 0.76, "BOOM", "execute / speculate", PANEL2, CYAN, CYAN, 14)
node(s, 2.68, 1.93, 1.72, 0.76, "GH_BUF", "packetize + seq", PANEL2, CYAN, CYAN, 13)
node(s, 4.82, 1.93, 1.80, 0.76, "GHM / CDC", "async queue", PANEL2, YELLOW, YELLOW, 13)
node(s, 7.02, 1.93, 1.58, 0.76, "Rocket ×4", "replay / compare", PANEL2, GREEN, GREEN, 13)
node(s, 9.02, 1.93, 1.72, 0.76, "Check result", "pass / reject", PANEL2, GREEN, GREEN, 13)
arrow(s,2.30,2.31,2.58,2.31,CYAN,1.8); arrow(s,4.42,2.31,4.72,2.31,CYAN,1.8); arrow(s,6.65,2.31,6.92,2.31,YELLOW,1.8); arrow(s,8.68,2.31,8.92,2.31,GREEN,1.8)

# gate / state blocks
box(s, 0.80, 3.20, 3.70, 2.62, PANEL, GRID, 0.08, 0.8)
text_box(s, 1.02, 3.42, 3.2, 0.25, "1 · Gate manager", 13, RED, True, margin=0)
rich_lines(s, 1.02, 3.85, 3.20, 1.60, [("MMIO / uncacheable A", RED, True), "校验通过前不允许 Put / Atomic", "记录请求顺序与 commit token", "设备接收 = 不可逆边界"], 10.5, WHITE, bullet=True, margin=0)

box(s, 4.82, 3.20, 3.70, 2.62, PANEL, GRID, 0.08, 0.8)
text_box(s, 5.04, 3.42, 3.2, 0.25, "2 · Versioned L1D", 13, ORANGE, True, margin=0)
rich_lines(s, 5.04, 3.85, 3.20, 1.60, [("line timestamp / epoch", ORANGE, True), "未确认 dirty bytes 不得 Release / Probe", "old value + byte mask → undo log", "必要时暂停主核并缩短 checkpoint"], 10.5, WHITE, bullet=True, margin=0)

box(s, 8.84, 3.20, 3.68, 2.62, PANEL, GRID, 0.08, 0.8)
text_box(s, 9.06, 3.42, 3.2, 0.25, "3 · Recovery controller", 13, CYAN, True, margin=0)
rich_lines(s, 9.06, 3.85, 3.18, 1.60, [("checkpoint", CYAN, True), "flush ROB / TLB / internal queues", "restore CSR + register + reservation", "replay from last consistent epoch"], 10.5, WHITE, bullet=True, margin=0)

# feedback arrows and outcome
arrow(s, 9.90, 2.72, 9.90, 3.12, GREEN, 1.7)
arrow(s, 4.62, 4.47, 4.62, 3.12, RED, 1.7)
text_box(s, 9.34, 2.82, 2.75, 0.22, "reject → rollback", 9.5, RED, True, align=PP_ALIGN.CENTER, margin=0)
text_box(s, 1.05, 6.24, 11.3, 0.25, "成功路径：execute → validate → release      错误路径：detect → stop propagation → undo / checkpoint restore → replay", 10.5, WHITE, True, align=PP_ALIGN.CENTER, font=MONO, margin=0)
footer(s, "08 / SYSTEM FRAMEWORK")

# Slide 9
s = add_slide()
title(s, 9, "Takeaways", "Barrier + Checkpoint + Undo Log：三条设计原则", "先定位传播路径，再决定屏障；先定义可恢复性，再决定允许多少投机。", "研究内容汇总；Method/EC.typ；Method/boom_dataflow_analysis.typ")

principles = [
    ("01", "路径优先", "先分析状态传播路径，再决定 Barrier 位置。", "dmem.req.fire → L1D → C → L2 → AXI / device", CYAN),
    ("02", "边界必检", "所有可能传播到不可恢复区域的状态修改必须经过校验。", "uncacheable Store / AMO：Barrier before commit", RED),
    ("03", "提前必撤销", "所有允许提前执行的修改必须具备 Undo 能力。", "checkpoint + old value + epoch / version", ORANGE),
]
for i, (num, head, desc, code, col) in enumerate(principles):
    yy = 1.98 + i*1.28
    box(s, 0.86, yy, 11.62, 0.98, PANEL, col, 0.08, 1.0)
    tag(s, 1.10, yy+0.31, 0.66, num, col, BG, 10)
    text_box(s, 2.00, yy+0.20, 1.45, 0.28, head, 15, col, True, margin=0)
    text_box(s, 3.60, yy+0.16, 7.75, 0.30, desc, 12, WHITE, True, margin=0)
    text_box(s, 3.60, yy+0.55, 7.95, 0.22, code, 9.5, MUTED, font=MONO, margin=0)

box(s, 0.86, 5.98, 11.62, 0.65, CYAN_D, CYAN, 0.08, 0.8)
text_box(s, 1.15, 6.17, 11.05, 0.23, "最终目标：BOOM 保持高吞吐，Rocket 提供独立确认；错误不越过不可逆边界，恢复回到一致状态。", 12, WHITE, True, align=PP_ALIGN.CENTER, margin=0)
text_box(s, 0.86, 6.76, 11.62, 0.22, "下一步验证：MMIO barrier · L1D timestamp/eviction gate · undo log/ECC · checkpoint replay · fault injection", 9.2, MUTED, font=MONO, align=PP_ALIGN.CENTER, margin=0)
footer(s, "09 / TAKEAWAYS")

prs.save(OUT)
print(OUT)
